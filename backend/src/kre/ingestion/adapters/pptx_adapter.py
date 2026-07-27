from pathlib import Path

from pptx import Presentation

from kre.models import Chunk


def parse(path: Path, document_id: str) -> list[Chunk]:
    presentation = Presentation(path)
    chunks: list[Chunk] = []
    for slide_number, slide in enumerate(presentation.slides, 1):
        for shape_number, shape in enumerate(slide.shapes):
            text = getattr(shape, "text", "").strip()
            if text:
                chunks.append(Chunk(
                    id=f"{document_id}:slide:{slide_number}:shape:{shape_number}",
                    document_id=document_id, source_format="pptx", text=text,
                    element_type="paragraph", page_number=slide_number,
                    location_reference=f"Slide: {slide_number}, Shape: {shape_number}",
                ))
        if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                chunks.append(Chunk(
                    id=f"{document_id}:slide:{slide_number}:notes",
                    document_id=document_id, source_format="pptx", text=notes,
                    element_type="caption", page_number=slide_number,
                    location_reference=f"Slide: {slide_number}, Speaker notes",
                ))
    return chunks
