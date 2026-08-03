from pathlib import Path
import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation

from kre.ingestion_lambda.format_router import route
from kre.ingestion_lambda.adapters.docx_adapter import parse as parse_docx
from kre.ingestion_lambda.adapters.pdf_adapter import parse as parse_pdf
from kre.ingestion_lambda.adapters.pptx_adapter import parse as parse_pptx
from kre.ingestion_lambda.adapters.xlsx_adapter import parse as parse_xlsx
from kre.ingestion_lambda.adapters.csv_adapter import parse as parse_csv
from kre.ingestion_lambda.parse_service import parse_file
from kre.ingestion_lambda.page_index_service import rank, score
from kre.shared.models import Chunk

DATA_DIR = Path(__file__).parent / "data"


def test_format_router_rejects_unsupported(tmp_path: Path):
    try:
        route(tmp_path / "file.unknown")
    except ValueError as error:
        assert "Unsupported format" in str(error)
    else:
        raise AssertionError("unsupported format was accepted")


def test_docx_heading_and_paragraph(tmp_path: Path):
    path = tmp_path / "sample.docx"
    document = DocxDocument()
    document.add_heading("Policy", level=1)
    document.add_paragraph("Refunds are available within thirty days.")
    document.save(path)
    chunks = parse_docx(path, "doc")
    assert [chunk.element_type for chunk in chunks] == ["heading", "paragraph"]
    assert chunks[1].location_reference == "Paragraph: 2"


def test_xlsx_uses_computed_values(tmp_path: Path):
    path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet["A1"] = 12
    sheet["B1"] = "=A1*2"
    workbook.calculation.fullCalcOnLoad = True
    workbook.save(path)
    chunks = parse_xlsx(path, "doc")
    assert all("=" not in chunk.text for chunk in chunks)


def test_pptx_speaker_notes_are_caption(tmp_path: Path):
    path = tmp_path / "sample.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[1]).notes_slide.notes_text_frame.text = "Presenter note"
    presentation.save(path)
    chunks = parse_pptx(path, "doc")
    assert any(chunk.element_type == "caption" for chunk in chunks)


def test_pageindex_prioritizes_heading_over_footnote():
    heading = Chunk("h", "d", "pdf", "refund policy", "heading")
    footnote = Chunk("f", "d", "pdf", "refund policy", "footnote")
    assert score(heading, "refund policy") / score(footnote, "refund policy") >= 2
    assert rank([footnote, heading], "refund policy")[0] == heading


def test_real_pdf_ingestion_data():
    pdf_path = DATA_DIR / "hdfc.pdf"
    if not pdf_path.exists():
        pytest.skip("hdfc.pdf not found in tests/data")
    try:
        document = parse_file(pdf_path, "pdf-doc-1")
    except RuntimeError as e:
        if "opendataloader_pdf is not installed" in str(e):
            pytest.skip("opendataloader-pdf not installed, skipping test")
        raise
    assert document.source_format == "pdf"
    assert len(document.chunks) > 50
    # Bounding box non-null check for PDF chunks per Phase 1 exit criteria
    sample_chunk = document.chunks[0]
    assert sample_chunk.bounding_box is not None
    assert "x1" in sample_chunk.bounding_box
    assert sample_chunk.page_number is not None


def test_real_docx_ingestion_data():
    docx_path = DATA_DIR / "Workflow Documentation.docx"
    if not docx_path.exists():
        pytest.skip("Workflow Documentation.docx not found in tests/data")
    document = parse_file(docx_path, "docx-doc-1")
    assert document.source_format == "docx"
    assert len(document.chunks) > 0
    assert document.chunks[0].location_reference.startswith("Paragraph:")


def test_real_pptx_ingestion_data():
    pptx_path = DATA_DIR / "submission.pptx"
    if not pptx_path.exists():
        pytest.skip("submission.pptx not found in tests/data")
    document = parse_file(pptx_path, "pptx-doc-1")
    assert document.source_format == "pptx"
    assert len(document.chunks) > 0
    assert "Slide:" in document.chunks[0].location_reference


def test_real_csv_ingestion_data():
    csv_paths = list(DATA_DIR.glob("*.csv"))
    if not csv_paths:
        pytest.skip("No CSV files found in tests/data")
    for csv_path in csv_paths:
        document = parse_file(csv_path, f"csv-{csv_path.stem}")
        assert document.source_format == "csv"
        assert len(document.chunks) > 0


def test_real_xlsx_ingestion_data():
    xlsx_paths = list(DATA_DIR.glob("*.xlsx"))
    if not xlsx_paths:
        pytest.skip("No XLSX file available in tests/data (skipped per user request)")
    for xlsx_path in xlsx_paths:
        document = parse_file(xlsx_path, f"xlsx-{xlsx_path.stem}")
        assert document.source_format == "xlsx"
        assert len(document.chunks) > 0
